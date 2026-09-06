"""Exportacao de um bloco (grafico + tabela de variacao) para PowerPoint:
slide 1 reproduz o layout da tela (grafico a esquerda, tabela de
variacao a direita); se a tabela nao couber inteira no slide 1 (Marca/
Submarca/Variante podem chegar a 30 categorias no top N), o restante
continua em slides seguintes, so com a tabela (largura cheia).

Reaproveita `compute_variations` (mesma funcao usada pela tabela HTML do
app) pra garantir que os numeros exportados sejam identicos aos exibidos
na tela - so muda a camada de apresentacao (python-pptx em vez de Dash).
A tabela vira uma tabela PPTX nativa (editavel no PowerPoint), estilizada
o mais proximo possivel da tela: mesmas cores de alta/baixa, valor
nominal em cinza entre parenteses, categoria em negrito.
"""

from __future__ import annotations

from io import BytesIO

import plotly.graph_objects as go
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from charts import YEARS_DEFAULT, compute_variations

_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)
_MARGIN = Inches(0.35)
_GAP = Inches(0.3)

# fracao da largura util (descontada a margem) reservada pro grafico -
# espelha a proporcao "flex: 5" grafico / "flex: 4" tabela do layout na
# tela (5/9 =~ 0.556)
_CHART_WIDTH_FRACTION = 5 / 9

# menos linhas por slide quando a tabela divide espaco com o grafico
# (coluna bem mais estreita que a tabela em tela cheia); nos slides de
# continuacao (so tabela, largura inteira) cabem mais
_ROWS_PER_SIDE_SLIDE = 9
_ROWS_PER_FULL_SLIDE = 15

_POSITIVE_RGB = RGBColor(0x1E, 0x8E, 0x5A)
_NEGATIVE_RGB = RGBColor(0xC2, 0x3B, 0x3B)
_NOMINAL_RGB = RGBColor(0x33, 0x33, 0x33)
_MUTED_RGB = RGBColor(0xAA, 0xAA, 0xAA)
_HEADER_RGB = RGBColor(0x22, 0x22, 0x22)


def _add_variation_paragraph(paragraph, value, suffix, nominal_text, font_size):
    """Preenche um paragrafo com a variacao colorida (seta + %/pp) seguida
    do valor nominal em cinza escuro entre parenteses - mesma convencao de
    `app._variation_span`."""
    if value is None:
        run = paragraph.add_run()
        run.text = "–"
        run.font.color.rgb = _MUTED_RGB
        run.font.size = font_size
        return
    color = _POSITIVE_RGB if value >= 0 else _NEGATIVE_RGB
    icon = "▲" if value >= 0 else "▼"
    run = paragraph.add_run()
    run.text = f"{icon} {value:+.1f}{suffix}"
    run.font.color.rgb = color
    run.font.size = font_size
    if nominal_text is not None:
        run2 = paragraph.add_run()
        run2.text = f" ({nominal_text})"
        run2.font.color.rgb = _NOMINAL_RGB
        run2.font.size = font_size


def _fill_variation_cell(cell, pct, share_pp, nominal, share_value, value_decimals, show_share, font_size):
    tf = cell.text_frame
    tf.word_wrap = True
    nominal_text = f"{nominal:,.{value_decimals}f}" if nominal is not None else None
    _add_variation_paragraph(tf.paragraphs[0], pct, "%", nominal_text, font_size)
    if show_share:
        share_text = f"{share_value:.1f}%" if share_value is not None else None
        p = tf.add_paragraph()
        _add_variation_paragraph(p, share_pp, "pp", share_text, font_size)


_HEADER_FILL_RGB = RGBColor(0xF5, 0xF5, 0xF5)
_BODY_FILL_RGB = RGBColor(0xFF, 0xFF, 0xFF)


def _style_table_plain(table, n_rows, n_cols):
    """Remove o banding/tema colorido padrao do PowerPoint pra tabela (fundo
    branco, cabecalho cinza bem claro) - o estilo padrao (faixas azuis
    alternadas) nao existe na tabela HTML da tela."""
    table.first_row = False
    table.horz_banding = False
    for i in range(n_rows):
        for j in range(n_cols):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _HEADER_FILL_RGB if i == 0 else _BODY_FILL_RGB
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _fill_table(table, header, rows_data, value_decimals, show_share, header_size, body_size):
    for j, text in enumerate(header):
        cell = table.cell(0, j)
        cell.text = text
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = header_size
        run.font.color.rgb = _HEADER_RGB

    for i, (cat, cells) in enumerate(rows_data, start=1):
        cat_cell = table.cell(i, 0)
        cat_cell.text = cat
        cat_cell.text_frame.word_wrap = True
        cat_run = cat_cell.text_frame.paragraphs[0].runs[0]
        cat_run.font.bold = True
        cat_run.font.size = body_size
        for j, (pct, share_pp, nominal, share_value) in enumerate(cells, start=1):
            _fill_variation_cell(
                table.cell(i, j), pct, share_pp, nominal, share_value, value_decimals, show_share, body_size,
            )


def _col_widths(total_width, n_year_cols, cat_width):
    year_width = Emu(int((total_width - cat_width) / max(n_year_cols, 1)))
    return [cat_width] + [year_width] * n_year_cols


def _chart_area_size():
    """(width, height) da caixa reservada pro grafico no slide 1 - fixa,
    independente da figura (o grafico e renderizado nessa mesma proporcao
    - ver `build_pptx` - entao sempre preenche a caixa inteira, sem
    letterboxing nem distorcao)."""
    usable_width = _SLIDE_WIDTH - 2 * _MARGIN
    chart_area_width = Emu(int(usable_width * _CHART_WIDTH_FRACTION))
    area_height = _SLIDE_HEIGHT - 2 * _MARGIN
    return chart_area_width, area_height


def _picture_box(img_w_px, img_h_px, box_left, box_top, box_width, box_height):
    """(left, top, width, height) do maior retangulo que cabe em
    `box_width`x`box_height` mantendo a proporcao de `img_w_px`x`img_h_px`
    (contain), centralizado na caixa. Precisa ser "contain" (nao esticar
    pra preencher a caixa inteira): o grafico usa margens em PIXELS fixos
    (ver `alluvial_stack_chart`/`line_evolution_chart` - margin=dict(...)
    reservado pros rotulos de categoria/conectores a direita, pill de
    variacao em cima etc.) calibrados pra sua largura/altura originais
    (760 x fig.layout.height); forcar uma proporcao diferente na
    renderizacao (ex.: exportar num tamanho que combine com a caixa do
    slide) faz essas margens ocuparem uma fracao errada da imagem e
    rotulos/linhas ficam cortados ou fora do lugar."""
    aspect = img_w_px / img_h_px
    if box_width / box_height > aspect:
        height = box_height
        width = Emu(int(box_height * aspect))
    else:
        width = box_width
        height = Emu(int(box_width / aspect))
    left = Emu(int(box_left + (box_width - width) / 2))
    top = Emu(int(box_top + (box_height - height) / 2))
    return left, top, width, height


def _add_combo_slide(prs, img_bytes, img_w_px, img_h_px, header, rows_data, value_decimals, show_share):
    """Slide com o grafico a esquerda e a tabela (`rows_data`, ja limitada
    ao que cabe) a direita - reproduz o layout da tela."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    chart_area_width, area_height = _chart_area_size()
    left, top, width, height = _picture_box(img_w_px, img_h_px, _MARGIN, _MARGIN, chart_area_width, area_height)
    slide.shapes.add_picture(BytesIO(img_bytes), left, top, width=width, height=height)

    table_left = _MARGIN + chart_area_width + _GAP
    table_width = _SLIDE_WIDTH - _MARGIN - table_left
    n_rows = len(rows_data) + 1
    n_cols = len(header)
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, table_left, _MARGIN, table_width, area_height)
    table = graphic_frame.table
    _style_table_plain(table, n_rows, n_cols)
    for i, col_width in enumerate(_col_widths(table_width, n_cols - 1, Inches(1.35))):
        table.columns[i].width = col_width
    _fill_table(table, header, rows_data, value_decimals, show_share, Pt(10), Pt(9))
    return slide


def _add_table_slide(prs, header, rows_data, value_decimals, show_share):
    """Slide de continuacao (so tabela, largura cheia) pro que nao coube
    no slide 1 junto com o grafico."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    n_rows = len(rows_data) + 1
    n_cols = len(header)
    width = _SLIDE_WIDTH - 2 * _MARGIN
    height = _SLIDE_HEIGHT - 2 * _MARGIN
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, _MARGIN, _MARGIN, width, height)
    table = graphic_frame.table
    _style_table_plain(table, n_rows, n_cols)
    for i, col_width in enumerate(_col_widths(width, n_cols - 1, Inches(2.8))):
        table.columns[i].width = col_width
    _fill_table(table, header, rows_data, value_decimals, show_share, Pt(12), Pt(11))
    return slide


def build_pptx(
    fig: go.Figure,
    categories: list[str],
    values: dict[str, dict[str, float]],
    additive: bool,
    value_decimals: int,
    totals_override: dict[str, float] | None = None,
) -> bytes:
    prs = Presentation()
    prs.slide_width = _SLIDE_WIDTH
    prs.slide_height = _SLIDE_HEIGHT

    # renderiza na proporcao ORIGINAL do grafico (760 - largura padrao dos
    # dois construtores de grafico em charts.py, sempre usada aqui - x a
    # altura de verdade, preservada em fig.layout.height mesmo depois do
    # app zerar fig.layout.width pra autosize) - as margens dos graficos
    # sao em pixels fixos, calibradas pra essa proporcao especificamente
    # (rotulos de categoria/conectores a direita, pill de variacao em
    # cima etc.); exportar numa proporcao diferente (ex.: pra bater com a
    # caixa do slide) faz as margens ocuparem uma fracao errada da imagem
    # e rotulos/linhas saem cortados ou fora do lugar - por isso o
    # encaixe na caixa do slide e "contain" (ver `_picture_box`), nao
    # esticado
    img_w_px = 760
    img_h_px = int(fig.layout.height or 640)
    img_bytes = fig.to_image(format="png", width=img_w_px, height=img_h_px, scale=3)

    header: list[str] = []
    rows_data: list[tuple[str, list]] = []
    if categories:
        variations = compute_variations(values, categories, YEARS_DEFAULT, additive, totals_override)
        year_pairs = [(YEARS_DEFAULT[i][1:], YEARS_DEFAULT[i + 1][1:]) for i in range(len(YEARS_DEFAULT) - 1)]
        header = ["Categoria"] + [f"{y0}→{y1}" for y0, y1 in year_pairs]
        rows_data = [
            (
                cat,
                [
                    (
                        variations[cat]["pct"][i],
                        variations[cat]["share_pp"][i],
                        variations[cat]["nominal"][i],
                        variations[cat]["share_value"][i],
                    )
                    for i in range(len(year_pairs))
                ],
            )
            for cat in categories
        ]

    first_chunk = rows_data[:_ROWS_PER_SIDE_SLIDE]
    rest = rows_data[_ROWS_PER_SIDE_SLIDE:]
    _add_combo_slide(prs, img_bytes, img_w_px, img_h_px, header, first_chunk, value_decimals, additive)

    for start in range(0, len(rest), _ROWS_PER_FULL_SLIDE):
        chunk = rest[start : start + _ROWS_PER_FULL_SLIDE]
        _add_table_slide(prs, header, chunk, value_decimals, additive)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
